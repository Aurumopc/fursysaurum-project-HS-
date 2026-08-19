"""
샘플 테스트 스크립트: DESKS_CHANCE.pptx 에서 이미지/텍스트를 추출하여
fursys-store/images/FBD/ 에 저장하고 fursys-store/data/products.json 의
"FBD" 항목을 업데이트한다.

대상 슬라이드 (샘플):
  Slide 1 : 표지 -> 시리즈명 확인
  Slide 2 : 공간 납품사례 -> case_01.jpg
  Slide 3 : 제품 상세 본문(단면형 독립 데스크) -> main.jpg, feat_01.jpg, feat_02.jpg,
            컬러옵션 아이콘, 품번/사이즈 표

레이아웃 규칙 (실측 좌표 기반, cm 단위, 슬라이드 42.0 x 29.7cm):
  - 좌측 컬럼: left < 15.0
      * top < 23.5  : 좌측 대형 메인 컷 이미지
      * top >= 23.5 : 좌측 하단 컬러 옵션/아이콘 이미지
  - 우측 컬럼: left >= 15.0
      * top < 23.5  : 특장점 영역 (2자리 숫자 배지 "01","02"... 가 이미지쪽/텍스트쪽에
        한 쌍씩 존재. 배지 위치로 이미지-제목-설명을 매칭한다)
      * top >= 23.5 : 하단 패널 (품번+사이즈 텍스트, ※ 비고 텍스트)

python-pptx 로 도형 트리를 순회하고 Pillow 로 이미지 blob 을 표준 JPG 로 저장한다.
"""

import io
import json
import os
import re
from collections import defaultdict
from datetime import datetime

from pptx import Presentation
from pptx.util import Emu
from PIL import Image, UnidentifiedImageError

PPTX_PATH = r"C:\Users\ultraman\AppData\Local\Temp\!RAONK\_읽기전용_DESKS_CHANCE.pptx"
PROJECT_ROOT = r"C:\Users\ultraman\OneDrive\Desktop\fursys-store"
PRODUCT_CODE = "FBD"
IMAGES_DIR = os.path.join(PROJECT_ROOT, "images", PRODUCT_CODE)
DATA_PATH = os.path.join(PROJECT_ROOT, "data", "products.json")

LEFT_X_CM = 15.0     # 좌/우 컬럼 경계
FOOTER_Y_CM = 23.5   # 본문/하단 패널 경계

PICTURE_TYPE = 13  # MSO_SHAPE_TYPE.PICTURE


def cm(emu_value):
    return Emu(emu_value).cm if emu_value is not None else None


def shape_text(shape):
    if shape.has_text_frame:
        return shape.text_frame.text.strip()
    return ""


def save_picture(shape, out_path_no_ext):
    """도형의 이미지 blob 을 Pillow 로 열어 JPG 로 저장. 벡터(EMF/WMF) 등
    Pillow 가 열 수 없는 포맷은 원본 확장자 그대로 저장한다."""
    image = shape.image
    blob = image.blob
    os.makedirs(os.path.dirname(out_path_no_ext), exist_ok=True)
    try:
        im = Image.open(io.BytesIO(blob))
        if im.mode in ("RGBA", "P", "LA"):
            im = im.convert("RGB")
        out_path = out_path_no_ext + ".jpg"
        im.save(out_path, "JPEG", quality=92)
        return out_path
    except UnidentifiedImageError:
        out_path = out_path_no_ext + f".{image.ext}"
        with open(out_path, "wb") as f:
            f.write(blob)
        print(f"  [경고] Pillow 로 디코딩 불가({image.ext}) -> 원본 형식으로 저장: {out_path}")
        return out_path


def rel(path):
    return os.path.relpath(path, PROJECT_ROOT).replace("\\", "/")


def extract_slide1_series(slide):
    texts = [shape_text(s) for s in slide.shapes if shape_text(s)]
    series = texts[0] if texts else None
    print(f"[Slide 1] 표지 텍스트: {texts}")
    print(f"[Slide 1] 시리즈명: {series}")
    return series


def extract_slide2_case(slide):
    pics = [s for s in slide.shapes if s.shape_type == PICTURE_TYPE]
    pics.sort(key=lambda s: (cm(s.width) or 0) * (cm(s.height) or 0), reverse=True)
    case_images = []
    for idx, pic in enumerate(pics[:1], start=1):
        out = save_picture(pic, os.path.join(IMAGES_DIR, f"case_{idx:02d}"))
        case_images.append(rel(out))
        print(f"[Slide 2] 공간 납품사례 이미지 저장: {out}")
    return case_images


def extract_slide3_product(slide):
    shapes = list(slide.shapes)

    left_shapes = [s for s in shapes if cm(s.left) is not None and cm(s.left) < LEFT_X_CM]
    right_shapes = [s for s in shapes if cm(s.left) is not None and cm(s.left) >= LEFT_X_CM]

    # ---- 제목 정보 ----
    title_texts = sorted(
        [s for s in left_shapes if shape_text(s) and cm(s.top) < 10],
        key=lambda s: cm(s.top),
    )
    product_category_title = shape_text(title_texts[0]) if len(title_texts) > 0 else None
    product_subtitle = shape_text(title_texts[1]) if len(title_texts) > 1 else None
    series_on_page = shape_text(title_texts[2]) if len(title_texts) > 2 else None
    print(f"[Slide 3] 제목: {product_category_title} / 부제: {product_subtitle} / 시리즈: {series_on_page}")

    # ---- 좌측 대형 메인 컷 ----
    left_main_pics = [
        s for s in left_shapes if s.shape_type == PICTURE_TYPE and cm(s.top) < FOOTER_Y_CM
    ]
    left_main_pics.sort(key=lambda s: (cm(s.width) or 0) * (cm(s.height) or 0), reverse=True)
    main_path = None
    if left_main_pics:
        out = save_picture(left_main_pics[0], os.path.join(IMAGES_DIR, "main"))
        main_path = rel(out)
        print(f"[Slide 3] 메인 컷 저장: {out}")

    # ---- 좌측 하단 컬러 옵션 아이콘 ----
    left_color_pics = [
        s for s in left_shapes if s.shape_type == PICTURE_TYPE and cm(s.top) >= FOOTER_Y_CM
    ]
    color_option_images = []
    for idx, pic in enumerate(left_color_pics, start=1):
        out = save_picture(pic, os.path.join(IMAGES_DIR, f"color_{idx:02d}"))
        color_option_images.append(rel(out))
        print(f"[Slide 3] 컬러 옵션 아이콘 저장: {out}")

    # ---- 우측 상단: 특장점(번호 배지로 이미지<->제목/설명 매칭) ----
    right_top = [s for s in right_shapes if cm(s.top) < FOOTER_Y_CM]
    right_bottom = [s for s in right_shapes if cm(s.top) >= FOOTER_Y_CM]

    badge_re = re.compile(r"^\d{2}$")
    badges = [s for s in right_top if badge_re.match(shape_text(s))]
    badge_groups = defaultdict(list)
    for b in badges:
        badge_groups[shape_text(b)].append(b)

    pics_top = [s for s in right_top if s.shape_type == PICTURE_TYPE]
    texts_top = [s for s in right_top if shape_text(s) and not badge_re.match(shape_text(s))]

    def find_nearby_pic(badge):
        for p in pics_top:
            if (
                abs(cm(p.left) - (cm(badge.left) - 0.14)) < 0.5
                and abs(cm(p.top) - (cm(badge.top) - 0.24)) < 0.5
            ):
                return p
        return None

    features = []
    feat_counter = 0
    for no in sorted(badge_groups.keys()):
        blist = badge_groups[no]
        img_badge_pic = None
        title_badge = None
        for b in blist:
            pic = find_nearby_pic(b)
            if pic is not None:
                img_badge_pic = pic
            else:
                title_badge = b

        if img_badge_pic is None or title_badge is None:
            print(f"  [경고] 배지 '{no}' 매칭 실패 (image={img_badge_pic}, title_badge={title_badge})")
            continue

        title_tb = None
        for t in texts_top:
            if (
                abs(cm(t.top) - cm(title_badge.top)) < 0.2
                and cm(t.left) > cm(title_badge.left)
                and (cm(t.width) or 0) > 2
            ):
                title_tb = t
                break

        desc_tb = None
        if title_tb is not None:
            title_bottom = cm(title_tb.top) + cm(title_tb.height)
            best_d = None
            for t in texts_top:
                if t is title_tb:
                    continue
                if cm(t.left) >= cm(title_badge.left) and cm(t.top) >= title_bottom - 0.3:
                    d = cm(t.top) - title_bottom
                    if d >= -0.3 and (best_d is None or d < best_d):
                        best_d = d
                        desc_tb = t

        feat_counter += 1
        out = save_picture(img_badge_pic, os.path.join(IMAGES_DIR, f"feat_{feat_counter:02d}"))
        feature = {
            "no": no,
            "title": shape_text(title_tb) if title_tb else None,
            "description": shape_text(desc_tb) if desc_tb else None,
            "image": rel(out),
        }
        features.append(feature)
        print(f"[Slide 3] 특장점 {no} 저장: {out} / 제목='{feature['title']}'")

    features.sort(key=lambda f: f["no"])

    # ---- 우측 하단: 품번/사이즈 표 + 비고 ----
    spec_table = []
    notes = []
    code_size_re = re.compile(r"([A-Z]{2,}\d+)\s+(\d+)·(\d+)·(\d+)")
    for s in right_bottom:
        txt = shape_text(s)
        if not txt:
            continue
        if "·" in txt and re.search(r"[A-Z]{2,}\d", txt):
            for m in code_size_re.finditer(txt):
                spec_table.append(
                    {
                        "product_code": m.group(1),
                        "width_mm": int(m.group(2)),
                        "depth_mm": int(m.group(3)),
                        "height_mm": int(m.group(4)),
                    }
                )
        elif txt.startswith("※"):
            notes.append(txt)

    print(f"[Slide 3] 품번/사이즈 표: {spec_table}")
    print(f"[Slide 3] 비고: {notes}")

    return {
        "product_category_title": product_category_title,
        "product_subtitle": product_subtitle,
        "series": series_on_page,
        "main_image": main_path,
        "color_option_images": color_option_images,
        "features": features,
        "spec_table": spec_table,
        "notes": notes,
    }


def update_products_json(series, case_images, slide3_data):
    os.makedirs(os.path.dirname(DATA_PATH), exist_ok=True)
    if os.path.exists(DATA_PATH):
        with open(DATA_PATH, "r", encoding="utf-8") as f:
            db = json.load(f)
    else:
        db = {}

    entry = db.get(PRODUCT_CODE, {})
    entry.update(
        {
            "product_code": PRODUCT_CODE,
            "series": series or slide3_data["series"],
            "product_name": slide3_data["product_category_title"],
            "product_subtitle": slide3_data["product_subtitle"],
            "category": "책상",
            "images": {
                "case": case_images,
                "main": slide3_data["main_image"],
                "features": [f["image"] for f in slide3_data["features"]],
                "color_options": slide3_data["color_option_images"],
            },
            "features": slide3_data["features"],
            "color_options": {
                "images": slide3_data["color_option_images"],
                "materials": [],  # 샘플 슬라이드(3)에는 별도 색상/마감 텍스트가 없어 비워둠
            },
            "spec_table": slide3_data["spec_table"],
            "notes": slide3_data["notes"],
            "source": {
                "pptx": os.path.basename(PPTX_PATH),
                "slides_used": [1, 2, 3],
            },
            "updated_at": datetime.now().isoformat(timespec="seconds"),
        }
    )
    db[PRODUCT_CODE] = entry

    with open(DATA_PATH, "w", encoding="utf-8") as f:
        json.dump(db, f, ensure_ascii=False, indent=2)

    print(f"\n[저장 완료] {DATA_PATH}")
    return entry


def main():
    print(f"PPTX 로드: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)
    print(f"슬라이드 크기: {cm(prs.slide_width):.1f}cm x {cm(prs.slide_height):.1f}cm, 총 {len(prs.slides)} 슬라이드\n")

    os.makedirs(IMAGES_DIR, exist_ok=True)

    series = extract_slide1_series(prs.slides[0])
    print()
    case_images = extract_slide2_case(prs.slides[1])
    print()
    slide3_data = extract_slide3_product(prs.slides[2])
    print()

    entry = update_products_json(series, case_images, slide3_data)

    print("\n===== products.json['FBD'] 최종 내용 =====")
    print(json.dumps(entry, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
